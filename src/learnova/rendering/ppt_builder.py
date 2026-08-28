"""
Learnova Modern Themed PPTX Builder
Generates PPTX decks with dynamic visual layouts (Flowcharts, Tables, Metric Cards, Quizzes)
and embeds PowerPoint OpenXML slide transition animations.
"""

import io
import os
import re
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

from learnova.rendering import layout as L
from learnova.textutils import strip_inline_markdown
from learnova.rendering.theme_engine import (
    get_theme, auto_detect_theme, select_slide_layout, resolve_theme,
    ColorPalette, THEMES, LAYOUT_STYLES,
)

def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


_IMAGE_KEEP_ALL = os.getenv("LEARNOVA_IMAGE_KEEP_ALL", "").lower() in {"1", "true", "yes", "on"}


def _image_decision(orig: dict):
    """
    Decide whether a figure attached to a slide should appear in the PPTX, using
    the SAME policy as the web deck (``ai.image_policy``) so the two outputs never
    disagree. Returns ``(show: bool, caption: str)``. Only a confident DROP
    (logo / divider / bullet icon) is hidden.
    """
    img = (orig or {}).get("image") if isinstance(orig, dict) else None
    if not img or not img.get("bytes"):
        return False, ""
    try:
        from learnova.ai.image_policy import (
            ImageMeta, decide_image_action, meta_from_bytes,
        )

        ext = str(img.get("ext", "png")).lower().lstrip(".")
        slide_text = " ".join(str(x) for x in (orig.get("text"),))
        try:
            meta = meta_from_bytes(
                img["bytes"], ext=ext,
                ocr_text=str(img.get("description", "")), slide_text=slide_text,
            )
        except Exception:
            meta = ImageMeta(ext=ext, ocr_text=str(img.get("description", "")),
                             slide_text=slide_text)
        action = decide_image_action(meta)
        if action.action == "DROP" and not _IMAGE_KEEP_ALL:
            return False, ""
        return True, (action.caption or str(img.get("description", ""))[:200])
    except Exception:
        # Policy must never cost us a figure — when in doubt, show it.
        return True, str(img.get("description", ""))[:200]

def _pptx_family(slide, sp, theme, band) -> bool:
    """
    Draw the Deck Director's chosen visual family as native PowerPoint shapes so
    the PPTX matches the web deck instead of always falling back to bullets.
    Returns True if it drew something. Best-effort — any failure returns False
    and the caller renders the plain bullet list.
    """
    try:
        family = getattr(sp, "family", "") or ""
        data = getattr(sp, "data", {}) or {}
        if not data:
            return False
        L_ = band.left + 0.4
        W_ = band.width - 0.8
        top = band.top + 0.15
        h = band.height - 0.3

        def _box(x, y, w, hh, text, *, fill=None, bold=False, size=13, color=None, align=None):
            shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(hh))
            if fill is not None:
                shp.fill.solid(); shp.fill.fore_color.rgb = fill
                shp.line.color.rgb = theme.primary_rgb
            tf = shp.text_frame; tf.word_wrap = True
            tf.text = str(text)
            p = tf.paragraphs[0]
            p.font.size = Pt(size); p.font.bold = bold
            p.font.color.rgb = color or theme.text_rgb
            if align is not None:
                p.alignment = align
            return shp

        # ── process / cycle — numbered stacked steps ─────────────────────────
        if family in {"PROCESS_LINEAR", "PROCESS_CYCLIC", "WORKED_EXAMPLE"}:
            steps = data.get("steps") or data.get("stages") or [
                r.get("step") for r in (data.get("rows") or []) if r.get("step")
            ]
            steps = [s for s in steps if str(s).strip()][:8]
            if len(steps) < 2:
                return False
            reasons = {i: r.get("reason", "") for i, r in enumerate(data.get("rows") or [])}
            row_h = min(0.9, h / len(steps))
            for i, s in enumerate(steps):
                y = top + i * row_h
                _box(L_, y, 0.5, row_h - 0.08, str(i + 1), fill=theme.primary_rgb,
                     bold=True, size=14, color=theme.accent_rgb, align=PP_ALIGN.CENTER)
                label = str(s)
                if reasons.get(i):
                    label += f"   —  {reasons[i]}"
                _box(L_ + 0.65, y, W_ - 0.65, row_h - 0.08, label, fill=theme.bg_rgb, size=12)
            return True

        # ── pros / cons — two columns ───────────────────────────────────────
        if family == "COMPARE_VISUAL":
            pros = [str(x) for x in (data.get("pros") or [])][:5]
            cons = [str(x) for x in (data.get("cons") or [])][:5]
            if not (pros and cons):
                return False
            cw = W_ / 2 - 0.15
            _box(L_, top, cw, 0.4, "Advantages", bold=True, size=13, color=theme.primary_rgb)
            _box(L_ + cw + 0.3, top, cw, 0.4, "Trade-offs", bold=True, size=13, color=theme.primary_rgb)
            _box(L_, top + 0.45, cw, h - 0.5, "\n".join(f"•  {p}" for p in pros),
                 fill=theme.bg_rgb, size=11)
            _box(L_ + cw + 0.3, top + 0.45, cw, h - 0.5, "\n".join(f"•  {c}" for c in cons),
                 fill=theme.bg_rgb, size=11)
            return True

        # ── pyramid — stacked bands, widest at the base ─────────────────────
        if family == "HIERARCHY_NEST":
            levels = [str(x) for x in (data.get("levels") or [])][:5]
            if len(levels) < 3:
                return False
            row_h = min(0.8, h / len(levels))
            for i, lv in enumerate(levels):
                frac = 0.45 + (i / max(1, len(levels) - 1)) * 0.5
                w = W_ * frac
                _box(L_ + (W_ - w) / 2, top + i * row_h, w, row_h - 0.1, lv,
                     fill=theme.primary_rgb, color=theme.accent_rgb, bold=True,
                     size=12, align=PP_ALIGN.CENTER)
            return True

        # ── card grid ──────────────────────────────────────────────────────
        if family == "LIST_STRUCTURED":
            cards = data.get("cards") or []
            cards = [c for c in cards if str(c.get("body", c)).strip()][:4]
            if len(cards) < 3:
                return False
            cw = W_ / len(cards) - 0.2
            for i, c in enumerate(cards):
                x = L_ + i * (cw + 0.25)
                head = str(c.get("heading", "")) if isinstance(c, dict) else ""
                body = str(c.get("body", c)) if isinstance(c, dict) else str(c)
                txt = (head + "\n" + body) if head else body
                _box(x, top, cw, h * 0.8, txt, fill=theme.bg_rgb, size=11)
            return True

        # ── bar chart — horizontal bars ────────────────────────────────────
        if family in {"CHART_CATEGORICAL", "CHART_RANKING", "CHART_PART_TO_WHOLE"}:
            pts = []
            for p in (data.get("points") or []):
                try:
                    pts.append((str(p.get("label", "")), float(p.get("value", 0))))
                except (ValueError, TypeError):
                    pass
            pts = pts[:8]
            if len(pts) < 2:
                return False
            mx = max(v for _, v in pts) or 1
            row_h = min(0.55, h / len(pts))
            for i, (lbl, val) in enumerate(pts):
                y = top + i * row_h
                _box(L_, y, W_ * 0.28, row_h - 0.06, lbl, size=11, align=PP_ALIGN.RIGHT)
                bw = max(0.15, (W_ * 0.6) * val / mx)
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Inches(L_ + W_ * 0.3), Inches(y + 0.04),
                                             Inches(bw), Inches(row_h - 0.14))
                bar.fill.solid(); bar.fill.fore_color.rgb = theme.primary_rgb
                bar.line.fill.background()
                _box(L_ + W_ * 0.3 + bw + 0.05, y, W_ * 0.1, row_h - 0.06,
                     ("%g" % val), size=10, color=theme.text_rgb)
            return True

        # ── timeline — dated rows ──────────────────────────────────────────
        if family == "TIMELINE":
            events = [e for e in (data.get("events") or []) if e.get("title")][:6]
            if len(events) < 3:
                return False
            row_h = min(0.7, h / len(events))
            for i, e in enumerate(events):
                y = top + i * row_h
                _box(L_, y, 1.2, row_h - 0.08, str(e.get("date", "")), bold=True,
                     size=12, color=theme.primary_rgb)
                _box(L_ + 1.35, y, W_ - 1.35, row_h - 0.08, str(e.get("title", "")),
                     fill=theme.bg_rgb, size=11)
            return True

        # ── mind map — centre + chips ──────────────────────────────────────
        if family == "MIND_MAP":
            center = str(data.get("center", ""))
            branches = [str(b) for b in (data.get("branches") or [])][:8]
            if not center or len(branches) < 3:
                return False
            _box(L_ + W_ / 2 - 1.4, top, 2.8, 0.7, center, fill=theme.primary_rgb,
                 color=theme.accent_rgb, bold=True, size=14, align=PP_ALIGN.CENTER)
            per_row = 3
            cw = W_ / per_row - 0.2
            for i, br in enumerate(branches):
                r, c = divmod(i, per_row)
                _box(L_ + c * (cw + 0.25), top + 0.9 + r * 0.75, cw, 0.6, br,
                     fill=theme.bg_rgb, size=11, align=PP_ALIGN.CENTER)
            return True

        return False
    except Exception:
        return False


def _add_slide_transition(slide):
    """Inject OpenXML transition tag for smooth slide entrance."""
    try:
        slide_elem = slide._element
        transition = OxmlElement('p:transition')
        transition.set('spd', 'med')
        push = OxmlElement('p:push')
        push.set('dir', 'l')
        transition.append(push)
        slide_elem.append(transition)
    except Exception:
        pass

def _add_header_bar(slide, title_text: str, theme: ColorPalette):
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = theme.primary_rgb
    header_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(L.fit_font_size([title_text], 12.3, 0.85, max_pt=28, min_pt=15))
    p.font.bold = True
    p.font.color.rgb = theme.accent_rgb

def _add_takeaway_bar(slide, takeaway_text: str, theme: ColorPalette):
    if not takeaway_text:
        return
    band = L.takeaway_band()
    tkw_box = slide.shapes.add_textbox(
        Inches(band.left), Inches(band.top), Inches(band.width), Inches(band.height))
    tkw_box.fill.solid()
    tkw_box.fill.fore_color.rgb = theme.primary_rgb
    tkw_box.line.color.rgb = theme.accent_rgb
    tkw_box.line.width = Pt(2)

    ttk = tkw_box.text_frame
    ttk.word_wrap = True
    tp = ttk.paragraphs[0]
    tp.text = f"Key Takeaway: {takeaway_text}"
    tp.font.color.rgb = theme.text_rgb
    tp.font.size = Pt(L.fit_font_size(
        [tp.text], band.width - 0.3, band.height - 0.15, max_pt=15, min_pt=9))
    tp.font.bold = True

def _add_figure_slide(prs, theme: ColorPalette, title_text: str, image_bytes: bytes,
                      caption: str = "") -> bool:
    """
    Emit a dedicated slide for an extracted figure.

    Layouts other than MINIMAL_TEXT fill their content area with cards, tables
    or metric blocks, leaving nowhere to put an image without overlapping. So
    rather than dropping the figure, we give it its own slide directly after
    the slide it was anchored to — keeping it beside the content that
    discusses it, at a readable size.
    """
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme.bg_rgb
        _add_slide_transition(slide)
        _add_header_bar(slide, f"{title_text} — Figure", theme)

        # Fit the image inside the content band, preserving aspect ratio, then
        # centre it. Setting height alone let python-pptx derive the width from
        # the aspect ratio unchecked: a wide figure came out 12.15" starting at
        # x=2.4", i.e. 1.2" past the right edge of the slide.
        box_w = L.SLIDE_W - 2 * L.MARGIN_X
        box_h = (6.0 if caption else L.SLIDE_H - L.BOTTOM_MARGIN) - L.CONTENT_TOP - L.GAP
        pic = slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(0), Inches(0))
        native_w, native_h = Emu(pic.width).inches, Emu(pic.height).inches
        scale = min(box_w / native_w, box_h / native_h) if native_w and native_h else 1.0
        draw_w, draw_h = native_w * scale, native_h * scale
        pic.width, pic.height = Inches(draw_w), Inches(draw_h)
        pic.left = Inches(L.MARGIN_X + (box_w - draw_w) / 2)
        pic.top = Inches(L.CONTENT_TOP + (box_h - draw_h) / 2)

        if caption:
            cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9))
            ctf = cap_box.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.text = caption[:300]
            cp.font.size = Pt(13)
            cp.font.italic = True
            cp.font.color.rgb = theme.text_rgb
        return True
    except Exception:
        return False


# A leading "Label: body" pair, which source decks use constantly
# ("Definition: ...", "Purpose: ...", "Cash Flows: ..."). The label is the most
# useful thing on the card and was being discarded.
_LABELLED = re.compile(r"^\s*([A-Z][\w&/()\- ]{2,34}?)\s*[:–—-]\s+(\S.*)$", re.S)

# Schema words that describe the field rather than name the content.
_META_LABELS = {
    "label", "item", "point", "bullet", "text", "title", "heading",
    "description", "desc", "value", "field", "key", "name", "content",
}


def _split_card_label(text: str, fallback: str) -> tuple[str, str]:
    """
    Split a bullet into (card heading, card body).

    Cards used to be headed "PILLAR 1", "STEP 2" — numbering that tells a
    reader nothing. Where the text carries its own label we promote it; the
    generic heading is only a fallback.
    """
    # Normalise here as well as upstream: a leftover emphasis marker sits
    # between the label and its colon ("Strategic Growth:* Helps ...") and
    # stops the pattern matching, silently costing the card its real heading.
    clean = strip_inline_markdown(text)

    # Models sometimes echo the schema word instead of the content's own label,
    # giving "Label: Cash Flows - Estimated inflows". Peel that off and let the
    # real label behind it be found.
    for _ in range(2):
        match = _LABELLED.match(clean)
        if not match:
            break
        label, body = match.group(1).strip(), match.group(2).strip()
        if not body or len(label.split()) > 5:
            break
        if label.lower() in _META_LABELS:
            clean = body           # retry against what followed
            continue
        return label.upper(), body
    return fallback, clean


def _apply_theme_fonts(prs, theme: ColorPalette) -> None:
    """
    Stamp the theme's typefaces onto every run in the deck.

    python-pptx sets no font name by default, so slides render in Calibri
    regardless of the chosen pairing. Rather than thread the family through
    ~40 individual paragraph writes, apply it once at the end: bold or large
    runs get the display face, everything else gets the body face.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            frames = []
            if shape.has_text_frame:
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        frames.append(cell.text_frame)

            for frame in frames:
                for paragraph in frame.paragraphs:
                    targets = list(paragraph.runs) or []
                    for run in targets:
                        size = run.font.size or paragraph.font.size
                        is_display = bool(run.font.bold) or (
                            size is not None and size >= Pt(24)
                        )
                        run.font.name = (
                            theme.heading_font if is_display else theme.body_font
                        )


def _add_inline_quiz(slide, quiz: dict, theme: ColorPalette, has_takeaway: bool) -> None:
    """
    Draw a checkpoint question as a band at the foot of the slide.

    Keeping the question beside the material it tests reads better than a
    standalone quiz slide, and it stops the deck inflating with interruptions.
    The four options sit in one row so the whole band stays shallow.
    """
    band = L.quiz_band(has_takeaway)

    backdrop = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(band.left), Inches(band.top), Inches(band.width), Inches(band.height))
    backdrop.fill.solid()
    backdrop.fill.fore_color.rgb = theme.card_bg_rgb
    backdrop.line.color.rgb = theme.accent_rgb
    backdrop.line.width = Pt(2)
    backdrop.text_frame.text = ""

    question = (quiz.get("question") or "").strip()
    label = f"Q{quiz.get('index', 1)}. {question}"

    q_box = slide.shapes.add_textbox(
        Inches(band.left + 0.18), Inches(band.top + 0.1),
        Inches(band.width - 0.36), Inches(0.72))
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    q_para = q_tf.paragraphs[0]
    q_para.text = label
    q_para.font.bold = True
    q_para.font.color.rgb = theme.text_rgb
    q_para.font.size = Pt(L.fit_font_size([label], band.width - 0.4, 0.7,
                                          max_pt=14, min_pt=9))

    options = [str(o).strip() for o in (quiz.get("options") or []) if str(o).strip()][:4]
    if not options:
        return

    letters = "ABCD"
    row_top = band.top + 0.9
    row_h = band.height - 1.02
    cells = L.grid_cells(len(options),
                         L.Box(band.left + 0.18, row_top, band.width - 0.36, row_h),
                         gap=0.14, max_per_row=4)

    # One size across every option keeps the row visually even.
    opt_pt = min(
        L.fit_font_size([o], c.width - 0.2, c.height - 0.08, max_pt=12, min_pt=7)
        for o, c in zip(options, cells)
    )

    for i, (option, cell) in enumerate(zip(options, cells)):
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cell.left), Inches(cell.top), Inches(cell.width), Inches(cell.height))
        chip.fill.solid()
        chip.fill.fore_color.rgb = theme.bg_rgb
        chip.line.color.rgb = theme.primary_rgb
        chip.line.width = Pt(1)

        tf = chip.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        # Strip any leading "A)" the model already produced, so it is not doubled.
        clean = re.sub(r"^\s*[A-Da-d][).:]\s*", "", option)
        para.text = f"{letters[i]}. {clean}"
        para.font.size = Pt(opt_pt)
        para.font.color.rgb = theme.text_rgb


def build_pptx(slides_data: list[dict], topic_title: str = "Learnova Presentation",
               theme_id: str = "auto", theme_spec: dict | None = None,
               deck_plan=None) -> bytes:
    theme = resolve_theme(topic_title, theme_id, theme_spec)

    # Deck Director → per-slide speaker notes (the PPTX notes pane). Best-effort.
    if deck_plan is None:
        try:
            from learnova.rendering.deck_director import plan_deck

            deck_plan = plan_deck(slides_data)
        except Exception:
            deck_plan = None

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    # ── Title slide ───────────────────────────────────────────────────────────
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = theme.primary_rgb

    accent = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.33), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent_rgb
    accent.line.fill.background()

    tf1 = title_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.8)).text_frame
    tf1.text = topic_title.upper()
    tf1.paragraphs[0].font.size = Pt(48)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = theme.accent_rgb
    tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

    tf2 = title_slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1)).text_frame
    tf2.text = f"Generated by Learnova AI Engine • Theme: {theme.name}"
    tf2.paragraphs[0].font.size = Pt(20)
    tf2.paragraphs[0].font.color.rgb = theme.text_rgb
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ── Content slides ────────────────────────────────────────────────────────
    prev_style = None
    for idx, data in enumerate(slides_data):
        orig = data.get("original", {})
        imp = data.get("improved", {})
        if not isinstance(imp, dict):
            imp = {}

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme.bg_rgb
        _add_slide_transition(slide)

        layout_type = imp.get("layout_type", "MINIMAL_TEXT").upper()
        style_pattern = select_slide_layout(idx, layout_type, prev_style)
        prev_style = style_pattern

        title_text = imp.get("title", orig.get("title", "Presentation Slide"))
        takeaway_text = imp.get("takeaway", "").strip()
        inline_quiz = imp.get("inline_quiz") or None

        # Speaker notes pane (key point, per-click cues, "read exactly", timing)
        _sp = deck_plan.by_index(idx) if deck_plan is not None else None
        if _sp and getattr(_sp, "speaker_notes", ""):
            try:
                slide.notes_slide.notes_text_frame.text = _sp.speaker_notes
            except Exception:
                pass

        # Geometry is derived per slide: the body shrinks when a takeaway bar
        # or a checkpoint band is also being drawn, so nothing overlaps.
        band = L.content_band(has_takeaway=bool(takeaway_text),
                              has_quiz=bool(inline_quiz))

        _add_header_bar(slide, title_text, theme)

        # ── 1. TABLE LAYOUT ──────────────────────────────────────────────────
        if layout_type == "TABLE" and imp.get("table_headers") and imp.get("table_rows"):
            headers = imp.get("table_headers", [])
            rows = imp.get("table_rows", [])

            # Lead-in bullets sit above the grid; the table takes what is left,
            # so the two never overlap.
            lead = [str(b).strip() for b in (imp.get("bullets") or []) if str(b).strip()]
            table_top = band.top
            if lead:
                lead_h = min(band.height * 0.4,
                             L.block_height(lead, band.width - 0.3, 13.0) + 0.1)
                lead_box = slide.shapes.add_textbox(
                    Inches(band.left), Inches(band.top),
                    Inches(band.width), Inches(lead_h))
                ltf = lead_box.text_frame
                ltf.word_wrap = True
                lead_pt = L.fit_font_size(lead, band.width - 0.3, lead_h,
                                          max_pt=14, min_pt=9)
                for i, line in enumerate(lead):
                    lp = ltf.add_paragraph() if i else ltf.paragraphs[0]
                    lp.text = line
                    lp.font.size = Pt(lead_pt)
                    lp.font.color.rgb = theme.text_rgb
                table_top = band.top + lead_h + 0.12

            if headers and rows:
                num_rows = len(rows) + 1
                num_cols = len(headers)
                avail_h = max(0.6, band.top + band.height - table_top)
                tbl_shape = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(band.left), Inches(table_top),
                    Inches(band.width), Inches(min(avail_h, num_rows * 0.42 + 0.3)))
                tbl = tbl_shape.table

                # Format Header Row
                for c_idx, h_text in enumerate(headers):
                    cell = tbl.cell(0, c_idx)
                    cell.text = str(h_text)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = theme.primary_rgb
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = theme.accent_rgb
                        p.font.bold = True
                        p.font.size = Pt(14)

                # Format Data Rows
                for r_idx, row_vals in enumerate(rows):
                    for c_idx, val in enumerate(row_vals):
                        if c_idx < num_cols:
                            cell = tbl.cell(r_idx + 1, c_idx)
                            cell.text = str(val)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = theme.card_bg_rgb
                            for p in cell.text_frame.paragraphs:
                                p.font.size = Pt(12)
                                p.font.color.rgb = theme.text_rgb

        # ── 2. METRIC STAT CALLOUT LAYOUT ─────────────────────────────────────
        elif layout_type == "METRIC":
            metric_val = imp.get("metric_value", "100%")
            metric_lbl = imp.get("metric_label", title_text)
            metric_desc = imp.get("metric_desc", takeaway_text)

            m_w = min(9.33, band.width)
            m_h = min(4.0, band.height)
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(band.left + (band.width - m_w) / 2), Inches(band.top),
                Inches(m_w), Inches(m_h))
            box.fill.solid()
            box.fill.fore_color.rgb = theme.primary_rgb
            box.line.color.rgb = theme.accent_rgb
            box.line.width = Pt(4)

            mtf = box.text_frame
            mtf.word_wrap = True

            p1 = mtf.paragraphs[0]
            p1.text = metric_val
            p1.font.size = Pt(72)
            p1.font.bold = True
            p1.font.color.rgb = theme.accent_rgb
            p1.alignment = PP_ALIGN.CENTER

            p2 = mtf.add_paragraph()
            p2.text = metric_lbl.upper()
            p2.font.size = Pt(24)
            p2.font.bold = True
            p2.font.color.rgb = theme.subtext_rgb
            p2.alignment = PP_ALIGN.CENTER

            p3 = mtf.add_paragraph()
            p3.text = metric_desc
            p3.font.size = Pt(16)
            p3.font.color.rgb = theme.text_rgb
            p3.alignment = PP_ALIGN.CENTER

        # ── 3. INTERLEAVED QUIZ CHECKPOINT ────────────────────────────────────
        elif layout_type == "QUIZ":
            q_text = imp.get("question", "Checkpoint Question")
            options = imp.get("options", [])
            correct_opt = imp.get("correct", "A")

            qbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.2))
            qtf = qbox.text_frame
            qtf.word_wrap = True
            qp = qtf.paragraphs[0]
            qp.text = f"❓ Question: {q_text}"
            qp.font.size = Pt(20)
            qp.font.bold = True
            qp.font.color.rgb = theme.primary_rgb

            card_positions = [
                (Inches(0.5), Inches(2.7)),
                (Inches(6.8), Inches(2.7)),
                (Inches(0.5), Inches(4.3)),
                (Inches(6.8), Inches(4.3)),
            ]
            for idx, opt_text in enumerate(options[:4]):
                c_left, c_top = card_positions[idx]
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, c_top, Inches(6.0), Inches(1.3))
                card.fill.solid()
                card.fill.fore_color.rgb = theme.primary_rgb if str(opt_text).startswith(correct_opt) else theme.card_bg_rgb
                card.line.color.rgb = theme.primary_rgb
                card.line.width = Pt(2)

                ctf = card.text_frame
                ctf.word_wrap = True
                cp = ctf.paragraphs[0]
                cp.text = str(opt_text)
                cp.font.size = Pt(16)
                cp.font.bold = True
                cp.font.color.rgb = theme.accent_rgb if str(opt_text).startswith(correct_opt) else theme.text_rgb

        # ── 4. FLOWCHART / PROCESS STEP CARDS ──────────────────────────────────
        elif layout_type == "FLOWCHART":
            mermaid_code = imp.get("mermaid_code", "")
            bullets = imp.get("bullets", [])

            steps = bullets if bullets else ["Step 1: Initiation", "Step 2: Execution", "Step 3: Completion"]
            steps = [s for s in steps if str(s).strip()][:6]
            cells = L.grid_cells(len(steps), band, max_per_row=4)
            # One shared size so every card in the row reads evenly.
            step_pt = min(
                L.fit_font_size([s], c.width - 0.3, c.height - 0.9,
                                max_pt=L.MAX_CARD_PT, min_pt=L.MIN_CARD_PT)
                for s, c in zip(steps, cells)
            ) if steps else L.MAX_CARD_PT

            for idx, (step_text, cell) in enumerate(zip(steps, cells)):
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cell.left), Inches(cell.top),
                    Inches(cell.width), Inches(cell.height))
                card.fill.solid()
                card.fill.fore_color.rgb = theme.primary_rgb if idx % 2 == 0 else theme.card_bg_rgb
                card.line.color.rgb = theme.primary_rgb
                card.line.width = Pt(2)

                ctf = card.text_frame
                ctf.word_wrap = True

                heading, body = _split_card_label(step_text, f"STEP {idx + 1}")

                sp1 = ctf.paragraphs[0]
                sp1.text = heading
                sp1.font.size = Pt(14)
                sp1.font.bold = True
                sp1.font.color.rgb = theme.accent_rgb if idx % 2 == 0 else theme.primary_rgb

                sp2 = ctf.add_paragraph()
                sp2.text = body
                sp2.font.size = Pt(step_pt)
                sp2.font.color.rgb = theme.text_rgb

        # ── 5. CARD GRID LAYOUT ──────────────────────────────────────────────
        elif layout_type == "CARD_GRID":
            bullets = imp.get("bullets", [])
            # No placeholder cards: an empty grid used to render three boxes
            # literally reading "Pillar 1", "Pillar 2", "Pillar 3".
            items = [b for b in (bullets or []) if str(b).strip()][:6]
            cells = L.grid_cells(len(items), band, max_per_row=4)
            card_pt = min(
                L.fit_font_size([i], c.width - 0.3, c.height - 0.8,
                                max_pt=L.MAX_CARD_PT, min_pt=L.MIN_CARD_PT)
                for i, c in zip(items, cells)
            )

            for idx, (item_text, cell) in enumerate(zip(items, cells)):
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cell.left), Inches(cell.top),
                    Inches(cell.width), Inches(cell.height))
                card.fill.solid()
                card.fill.fore_color.rgb = theme.card_bg_rgb
                card.line.color.rgb = theme.primary_rgb
                card.line.width = Pt(2)

                ctf = card.text_frame
                ctf.word_wrap = True

                heading, body = _split_card_label(item_text, f"{idx + 1:02d}")

                sp1 = ctf.paragraphs[0]
                sp1.text = heading
                sp1.font.size = Pt(14)
                sp1.font.bold = True
                sp1.font.color.rgb = theme.primary_rgb

                sp2 = ctf.add_paragraph()
                sp2.text = body
                sp2.font.size = Pt(card_pt)
                sp2.font.color.rgb = theme.text_rgb

        # ── 5b. EXPANDED VISUAL FAMILY (from the Deck Director) ───────────────
        elif _sp and getattr(_sp, "family", None) and getattr(_sp, "data", None) and \
                _pptx_family(slide, _sp, theme, band):
            pass  # drew a native structured visual; bullets are represented in it

        # ── 6. DEFAULT / MINIMAL TEXT + IMAGE COLUMN LAYOUT ───────────────────
        else:
            _show_img, _img_caption = _image_decision(orig)
            has_img = _show_img
            text_area, image_area = L.split_text_image(band, has_img)

            bullets = [str(b).replace("▪", "").strip()
                       for b in (imp.get("bullets") or []) if str(b).strip()]

            body_box = slide.shapes.add_textbox(
                Inches(text_area.left), Inches(text_area.top),
                Inches(text_area.width), Inches(text_area.height))
            btf = body_box.text_frame
            btf.word_wrap = True

            # One size for the whole list, chosen so every bullet fits the band.
            body_pt = L.fit_font_size(bullets, text_area.width - 0.35,
                                      text_area.height - 0.1,
                                      max_pt=L.MAX_BODY_PT, min_pt=L.MIN_BODY_PT)

            for i, clean_b in enumerate(bullets):
                p = btf.add_paragraph() if i > 0 else btf.paragraphs[0]
                p.text = clean_b
                p.font.size = Pt(body_pt)
                p.font.color.rgb = theme.text_rgb
                p.space_after = Pt(max(4, body_pt * 0.45))
                pPr = p._p.get_or_add_pPr()
                SubElement(pPr, "a:buChar", char="•")

            if has_img:
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(orig["image"]["bytes"]),
                        Inches(image_area.left), Inches(image_area.top),
                        width=Inches(image_area.width))
                except Exception:
                    pass

        if inline_quiz:
            _add_inline_quiz(slide, inline_quiz, theme, bool(takeaway_text))

        _add_takeaway_bar(slide, takeaway_text, theme)

        # A figure anchored to a non-MINIMAL layout gets its own slide, placed
        # immediately after so it stays with the content it belongs to.
        image = orig.get("image") if isinstance(orig, dict) else None
        if layout_type != "MINIMAL_TEXT" and image and image.get("bytes"):
            _fig_show, _fig_caption = _image_decision(orig)
            if _fig_show:
                _add_figure_slide(
                    prs, theme, title_text, image["bytes"],
                    caption=(_fig_caption or image.get("description") or "")[:300],
                )

    # ── Closing slide ─────────────────────────────────────────────────────────
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    end_slide.background.fill.solid()
    end_slide.background.fill.fore_color.rgb = theme.primary_rgb
    _add_slide_transition(end_slide)

    accent2 = end_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.33), Inches(0.12))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = theme.accent_rgb

    tfe = end_slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5)).text_frame
    tfe.text = "THANK YOU"
    tfe.paragraphs[0].font.size = Pt(60)
    tfe.paragraphs[0].font.bold = True
    tfe.paragraphs[0].font.color.rgb = theme.accent_rgb
    tfe.paragraphs[0].alignment = PP_ALIGN.CENTER

    _apply_theme_fonts(prs, theme)

    # ── Optional per-shape entrance animations (LEARNOVA_PPTX_ANIM=1) ─────────
    try:
        from learnova.rendering.pptx_animation import (
            animations_enabled,
            apply_click_builds,
            strip_all_timing,
        )

        if animations_enabled():
            content_slides = list(prs.slides)[1:-1]  # skip title + thank-you
            for idx, slide in enumerate(content_slides):
                sp = deck_plan.by_index(idx) if deck_plan is not None else None
                n_groups = len((getattr(sp, "animation", {}) or {}).get("steps", []) or []) if sp else 0
                shape_ids = [s.shape_id for s in slide.shapes][2:8]
                if n_groups >= 2:
                    shape_ids = shape_ids[: max(2, n_groups)]
                apply_click_builds(slide, shape_ids)

            # Validate: the file must still round-trip through python-pptx.
            check = io.BytesIO()
            prs.save(check)
            check.seek(0)
            try:
                Presentation(check)
            except Exception:
                strip_all_timing(prs)
    except Exception:
        pass

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()